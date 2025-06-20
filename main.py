

# ✅ Import everything
import os
import getpass
from PyPDF2 import PdfReader
from pptx import Presentation
from langchain.text_splitter import CharacterTextSplitter
from langchain.embeddings import HuggingFaceEmbeddings
from langchain.vectorstores import FAISS
from langchain.llms import Together
from langchain.chains.question_answering import load_qa_chain
import gradio as gr

# ✅ Securely input Together.ai API key
os.environ["TOGETHER_API_KEY"] = getpass.getpass("Enter Together.ai API Key: ")

# ✅ PDF/PPTX Reader
def extract_text(file_path):
    if file_path.endswith(".pdf"):
        reader = PdfReader(file_path)
        text = ""
        for page in reader.pages:
            text += page.extract_text() + "\n"
        return text
    elif file_path.endswith(".pptx"):
        prs = Presentation(file_path)
        text = ""
        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    text += shape.text + "\n"
        return text
    else:
        return ""

# ✅ Chunk text for embeddings
def chunk_text(text):
    splitter = CharacterTextSplitter(
        separator="\n",
        chunk_size=500,
        chunk_overlap=100
    )
    return splitter.split_text(text)

# ✅ Create FAISS vector DB from chunks
def create_vectorstore(chunks):
    embeddings = HuggingFaceEmbeddings(model_name="all-MiniLM-L6-v2")
    return FAISS.from_texts(chunks, embedding=embeddings)

# ✅ QA chain using Together.ai model
def setup_qa_chain():
    llm = Together(
        model="mistralai/Mistral-7B-Instruct-v0.1",
        temperature=0.5,
        max_tokens=512
    )
    return load_qa_chain(llm, chain_type="stuff")

# ✅ Main logic for processing and answering
def qa_app(file, question):
    text = extract_text(file.name)
    if not text.strip():
        return "Unsupported file or empty content."

    chunks = chunk_text(text)
    vectorstore = create_vectorstore(chunks)
    docs = vectorstore.similarity_search(question)
    qa_chain = setup_qa_chain()
    result = qa_chain.run(input_documents=docs, question=question)
    return result

#✅ Gradio interface
gr.Interface(
    fn=qa_app,
    inputs=[
        gr.File(label="Upload PDF or PPTX"),
        gr.Textbox(label="Ask a Question")
    ],
    outputs="text",
    title="📚PDFPal -Chat with PDF/PPTX using Together.ai",
    description="Upload a file and ask questions. Powered by Mistral-7B on Together.ai"
).launch()

